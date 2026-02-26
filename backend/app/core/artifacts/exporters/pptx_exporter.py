"""PPTX exporter — converts DeckOutlineContent into polished presentations.

Uses python-pptx to programmatically build slides with professional
layouts, color scheme, and typography.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.models.artifact import ArtifactType

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------

_BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)  # Deep navy
_BG_ACCENT = RGBColor(0x16, 0x21, 0x3E)  # Slightly lighter navy
_TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)  # Soft white
_TEXT_MUTED = RGBColor(0xA0, 0xA4, 0xB8)  # Muted gray
_ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)  # Bright blue
_ACCENT_TEAL = RGBColor(0x00, 0xD2, 0xC6)  # Teal
_ACCENT_VIOLET = RGBColor(0x7C, 0x3A, 0xED)  # Violet

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout detection keywords
_TITLE_KEYWORDS = {"title", "cover", "intro", "introduction"}
_SECTION_KEYWORDS = {"agenda", "overview", "outline", "roadmap"}
_CLOSING_KEYWORDS = {"ask", "closing", "thank", "contact", "q&a", "next steps"}
_TWO_COL_KEYWORDS = {
    "comparison", "before", "after", "vs", "versus", "pros", "cons",
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _detect_layout(title: str, slide_index: int, total: int) -> str:
    """Determine which layout to use based on title and position."""
    t = title.lower().strip()
    if slide_index == 0:
        return "title"
    if slide_index == total - 1 and any(k in t for k in _CLOSING_KEYWORDS):
        return "closing"
    if any(k in t for k in _TITLE_KEYWORDS):
        return "section"
    if any(k in t for k in _SECTION_KEYWORDS):
        return "section"
    if any(k in t for k in _TWO_COL_KEYWORDS):
        return "two_column"
    if any(k in t for k in _CLOSING_KEYWORDS):
        return "closing"
    return "content"


def _set_slide_bg(slide: Any, color: RGBColor) -> None:
    """Set a solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(
    slide: Any,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: int = 18,
    color: RGBColor = _TEXT_WHITE,
    bold: bool = False,
    alignment: int = PP_ALIGN.LEFT,
) -> Any:
    """Add a text box to the slide and return the shape."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None  # type: ignore[assignment]

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txbox


def _add_accent_bar(
    slide: Any,
    left: int,
    top: int,
    width: int,
    height: int,
    color: RGBColor = _ACCENT_BLUE,
) -> None:
    """Add a thin accent rectangle."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _build_title_slide(prs: Presentation, title: str, points: list[str]) -> None:
    """Full-bleed title slide with company name and tagline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, _BG_DARK)

    # Accent bar at top
    _add_accent_bar(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), _ACCENT_BLUE)

    # Title
    _add_text_box(
        slide,
        Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
        title,
        font_size=44, bold=True, color=_TEXT_WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle / tagline from first key point
    if points:
        _add_text_box(
            slide,
            Inches(2), Inches(3.8), Inches(9), Inches(1),
            points[0],
            font_size=22, color=_TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )

    # Remaining points as small text
    if len(points) > 1:
        _add_text_box(
            slide,
            Inches(3), Inches(5.2), Inches(7), Inches(1.5),
            " · ".join(points[1:]),
            font_size=14, color=_TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )


def _build_content_slide(
    prs: Presentation, title: str, points: list[str],
) -> None:
    """Standard content slide with title + bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG_DARK)

    # Accent bar
    _add_accent_bar(
        slide, Inches(0.8), Inches(0.6),
        Inches(0.06), Inches(0.5), _ACCENT_BLUE,
    )

    # Title
    _add_text_box(
        slide,
        Inches(1.1), Inches(0.5), Inches(10), Inches(0.8),
        title,
        font_size=32, bold=True, color=_TEXT_WHITE,
    )

    # Divider line
    _add_accent_bar(
        slide, Inches(1.1), Inches(1.4),
        Inches(2), Inches(0.03), _ACCENT_TEAL,
    )

    # Bullet points
    y = Inches(1.8)
    for i, point in enumerate(points):
        accent = _ACCENT_BLUE if i % 2 == 0 else _ACCENT_TEAL
        # Bullet dot
        _add_accent_bar(
            slide, Inches(1.2), y + Pt(8),
            Pt(8), Pt(8), accent,
        )
        # Text
        _add_text_box(
            slide,
            Inches(1.6), y, Inches(10), Inches(0.6),
            point,
            font_size=18, color=_TEXT_WHITE,
        )
        y += Inches(0.65)


def _build_section_slide(
    prs: Presentation, title: str, points: list[str],
) -> None:
    """Section divider slide — large title, centered."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG_ACCENT)

    _add_accent_bar(slide, 0, 0, SLIDE_WIDTH, Inches(0.05), _ACCENT_VIOLET)

    _add_text_box(
        slide,
        Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
        title,
        font_size=40, bold=True, color=_TEXT_WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    if points:
        _add_text_box(
            slide,
            Inches(2), Inches(4.0), Inches(9), Inches(1.5),
            "\n".join(points[:3]),
            font_size=18, color=_TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )


def _build_two_column_slide(
    prs: Presentation, title: str, points: list[str],
) -> None:
    """Two-column layout — splits points evenly."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG_DARK)

    _add_accent_bar(
        slide, Inches(0.8), Inches(0.6),
        Inches(0.06), Inches(0.5), _ACCENT_VIOLET,
    )

    _add_text_box(
        slide,
        Inches(1.1), Inches(0.5), Inches(10), Inches(0.8),
        title,
        font_size=32, bold=True, color=_TEXT_WHITE,
    )

    _add_accent_bar(
        slide, Inches(1.1), Inches(1.4),
        Inches(2), Inches(0.03), _ACCENT_TEAL,
    )

    mid = max(len(points) // 2, 1)
    left_pts = points[:mid]
    right_pts = points[mid:]

    # Left column
    y = Inches(1.8)
    for pt in left_pts:
        _add_text_box(
            slide,
            Inches(1.2), y, Inches(5), Inches(0.55),
            f"→  {pt}",
            font_size=16, color=_TEXT_WHITE,
        )
        y += Inches(0.6)

    # Right column
    y = Inches(1.8)
    for pt in right_pts:
        _add_text_box(
            slide,
            Inches(7), y, Inches(5), Inches(0.55),
            f"→  {pt}",
            font_size=16, color=_TEXT_WHITE,
        )
        y += Inches(0.6)


def _build_closing_slide(
    prs: Presentation, title: str, points: list[str],
) -> None:
    """Closing / Ask slide — prominent CTA."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG_ACCENT)

    _add_accent_bar(slide, 0, Inches(7.3), SLIDE_WIDTH, Inches(0.2), _ACCENT_BLUE)

    _add_text_box(
        slide,
        Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
        title,
        font_size=40, bold=True, color=_TEXT_WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    if points:
        _add_text_box(
            slide,
            Inches(2), Inches(3.5), Inches(9), Inches(2.5),
            "\n".join(points),
            font_size=20, color=_TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )


_LAYOUT_BUILDERS = {
    "title": _build_title_slide,
    "content": _build_content_slide,
    "section": _build_section_slide,
    "two_column": _build_two_column_slide,
    "closing": _build_closing_slide,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


class PptxExporter:
    """Export DeckOutlineContent to .pptx bytes."""

    def export(
        self,
        artifact_type: ArtifactType,
        title: str,
        content: dict[str, Any],
    ) -> bytes:
        """Generate PPTX bytes from artifact content.

        Only supports ``ArtifactType.DECK_OUTLINE``.
        """
        if artifact_type != ArtifactType.DECK_OUTLINE:
            msg = (
                f"PptxExporter only supports deck_outline, "
                f"got {artifact_type.value}"
            )
            raise ValueError(msg)

        slides_data: list[dict[str, Any]] = content.get("slides", [])
        if not slides_data:
            msg = "No slides in deck_outline content"
            raise ValueError(msg)

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        total = len(slides_data)
        for i, slide_info in enumerate(slides_data):
            slide_title = slide_info.get("title", f"Slide {i + 1}")
            key_points: list[str] = slide_info.get("key_points", [])
            speaker_notes: str = slide_info.get("speaker_notes", "")

            layout = _detect_layout(slide_title, i, total)
            builder = _LAYOUT_BUILDERS[layout]
            builder(prs, slide_title, key_points)

            # Attach speaker notes
            if speaker_notes:
                last_slide = prs.slides[-1]
                notes_slide = last_slide.notes_slide
                notes_slide.notes_text_frame.text = speaker_notes

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()


pptx_exporter = PptxExporter()
