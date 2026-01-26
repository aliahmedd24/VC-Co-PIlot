"""Text extraction from various document formats."""

import csv
from io import BytesIO


def extract_text(content: bytes, mime_type: str) -> str:
    """Extract text from document content based on MIME type.

    Args:
        content: Raw file content as bytes.
        mime_type: MIME type of the file.

    Returns:
        Extracted text content.

    Raises:
        ValueError: If MIME type is not supported.
    """
    extractors = {
        "application/pdf": extract_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_xlsx,
        "application/vnd.ms-excel": extract_xlsx,
        "text/csv": extract_csv,
        "text/plain": lambda c: c.decode("utf-8", errors="ignore"),
    }

    extractor = extractors.get(mime_type)
    if not extractor:
        raise ValueError(f"Unsupported MIME type: {mime_type}")

    return extractor(content)


def extract_pdf(content: bytes) -> str:
    """Extract text from PDF content.

    Args:
        content: PDF file as bytes.

    Returns:
        Extracted text.
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError("pypdf not installed. Run: pip install pypdf") from None

    reader = PdfReader(BytesIO(content))
    text_parts = []

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)

    return "\n\n".join(text_parts)


def extract_docx(content: bytes) -> str:
    """Extract text from DOCX content.

    Args:
        content: DOCX file as bytes.

    Returns:
        Extracted text.
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx not installed. Run: pip install python-docx"
        ) from None

    doc = Document(BytesIO(content))
    text_parts = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text.strip())

    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)

    return "\n\n".join(text_parts)


def extract_pptx(content: bytes) -> str:
    """Extract text from PowerPoint content.

    Args:
        content: PPTX file as bytes.

    Returns:
        Extracted text.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx not installed. Run: pip install python-pptx"
        ) from None

    prs = Presentation(BytesIO(content))
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        if len(slide_texts) > 1:  # Has content beyond just slide number
            text_parts.append("\n".join(slide_texts))

    return "\n\n".join(text_parts)


def extract_xlsx(content: bytes) -> str:
    """Extract text from Excel content.

    Args:
        content: XLSX file as bytes.

    Returns:
        Extracted text.
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl") from None

    wb = openpyxl.load_workbook(BytesIO(content), data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_texts = [f"[Sheet: {sheet_name}]"]

        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            row_text = " | ".join(v for v in row_values if v)
            if row_text:
                sheet_texts.append(row_text)

        if len(sheet_texts) > 1:
            text_parts.append("\n".join(sheet_texts))

    return "\n\n".join(text_parts)


def extract_csv(content: bytes) -> str:
    """Extract text from CSV content.

    Args:
        content: CSV file as bytes.

    Returns:
        Extracted text.
    """
    text = content.decode("utf-8", errors="ignore")
    reader = csv.reader(text.splitlines())

    rows = []
    for row in reader:
        row_text = " | ".join(cell.strip() for cell in row if cell.strip())
        if row_text:
            rows.append(row_text)

    return "\n".join(rows)


# Supported MIME types
SUPPORTED_MIME_TYPES = [
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "text/csv",
    "text/plain",
]


def is_supported_mime_type(mime_type: str) -> bool:
    """Check if a MIME type is supported for extraction."""
    return mime_type in SUPPORTED_MIME_TYPES
