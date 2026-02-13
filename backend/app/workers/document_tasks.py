import io
import re
import uuid

import structlog
from sqlalchemy import select

from app.workers.celery_app import celery_app

logger = structlog.get_logger()

CHUNK_SIZE = 500  # tokens (approximate by words)
CHUNK_OVERLAP = 50


def _extract_text_pdf(file_data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _extract_text_docx(file_data: bytes) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(file_data))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_text_pptx(file_data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_data))
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
    return "\n\n".join(text_parts)


def _extract_text_xlsx(file_data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append("\t".join(cells))
    return "\n".join(rows)


def _extract_text_plain(file_data: bytes) -> str:
    return file_data.decode("utf-8", errors="replace")


EXTRACTORS: dict[str, object] = {
    "application/pdf": _extract_text_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_text_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_text_pptx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_text_xlsx,
    "text/plain": _extract_text_plain,
    "text/csv": _extract_text_plain,
}


def _chunk_text(text: str) -> list[str]:
    paragraphs = re.split(r"\n\s*\n", text)
    chunks: list[str] = []
    current_chunk: list[str] = []
    current_size = 0

    for paragraph in paragraphs:
        words = paragraph.split()
        if not words:
            continue

        # If a single paragraph exceeds chunk size, split it into sub-chunks
        if len(words) > CHUNK_SIZE:
            # Flush current chunk first
            if current_chunk:
                chunks.append(" ".join(current_chunk))
                overlap_words = (
                    current_chunk[-CHUNK_OVERLAP:]
                    if len(current_chunk) > CHUNK_OVERLAP
                    else current_chunk[:]
                )
                current_chunk = overlap_words
                current_size = len(current_chunk)

            # Split large paragraph
            for i in range(0, len(words), CHUNK_SIZE - CHUNK_OVERLAP):
                sub = words[i : i + CHUNK_SIZE]
                if current_chunk:
                    current_chunk.extend(sub)
                    chunks.append(" ".join(current_chunk))
                    current_chunk = sub[-CHUNK_OVERLAP:]
                    current_size = len(current_chunk)
                else:
                    chunks.append(" ".join(sub))
                    current_chunk = sub[-CHUNK_OVERLAP:]
                    current_size = len(current_chunk)
            continue

        word_count = len(words)

        if current_size + word_count > CHUNK_SIZE and current_chunk:
            chunks.append(" ".join(current_chunk))

            overlap_words = (
                current_chunk[-CHUNK_OVERLAP:]
                if len(current_chunk) > CHUNK_OVERLAP
                else current_chunk[:]
            )
            current_chunk = overlap_words
            current_size = len(current_chunk)

        current_chunk.extend(words)
        current_size += word_count

    if current_chunk and (not chunks or len(current_chunk) > CHUNK_OVERLAP):
        chunks.append(" ".join(current_chunk))

    return chunks


@celery_app.task(name="process_document", bind=True, max_retries=3)  # type: ignore[untyped-decorator]
def process_document(self: object, document_id: str) -> dict[str, str]:
    from sqlalchemy import create_engine
    from sqlalchemy.orm import sessionmaker

    from app.config import settings
    from app.models.document import Document, DocumentChunk, DocumentStatus
    from app.models.venture import Venture
    from app.services.embedding_service import embedding_service
    from app.services.storage_service import storage_service

    sync_url = (
        settings.database_url
        .replace("+asyncpg", "+psycopg2")
        .replace("postgresql+psycopg2", "postgresql")
    )
    engine = create_engine(sync_url)
    session_factory = sessionmaker(bind=engine)

    with session_factory() as db:
        doc = db.execute(
            select(Document).where(Document.id == uuid.UUID(document_id))
        ).scalar_one_or_none()
        if doc is None:
            logger.error("document_not_found", document_id=document_id)
            return {"status": "error", "message": "Document not found"}

        try:
            # Set processing status
            doc.status = DocumentStatus.PROCESSING
            db.commit()

            # Download from S3
            file_data = storage_service.download_file(doc.storage_key)

            # Extract text
            extractor = EXTRACTORS.get(doc.mime_type)
            if extractor is None:
                raise ValueError(f"No extractor for mime type: {doc.mime_type}")
            text = extractor(file_data)  # type: ignore[operator]

            if not text.strip():
                raise ValueError("No text extracted from document")

            # Chunk text
            chunks = _chunk_text(text)
            logger.info("text_chunked", document_id=document_id, chunk_count=len(chunks))

            # Get venture_id for this workspace
            venture = db.execute(
                select(Venture).where(Venture.workspace_id == doc.workspace_id)
            ).scalar_one_or_none()

            if venture is None:
                raise ValueError("No venture found for workspace")

            # Embed chunks
            embeddings = embedding_service.embed_batch(chunks)

            # Insert chunks
            for idx, (chunk_text, embedding) in enumerate(zip(chunks, embeddings, strict=True)):
                chunk = DocumentChunk(
                    document_id=doc.id,
                    venture_id=venture.id,
                    content=chunk_text,
                    embedding=embedding,
                    chunk_index=idx,
                    chunk_metadata={"source": doc.name, "chunk_index": idx},
                )
                db.add(chunk)

            doc.status = DocumentStatus.INDEXED
            db.commit()

            logger.info(
                "document_processed",
                document_id=document_id,
                chunks_created=len(chunks),
            )

            # Chain entity extraction task
            try:
                extract_entities.delay(document_id)
                logger.info("entity_extraction_enqueued", document_id=document_id)
            except Exception as chain_exc:
                logger.warning(
                    "entity_extraction_enqueue_failed",
                    document_id=document_id,
                    error=str(chain_exc),
                )

            return {"status": "success", "chunks": str(len(chunks))}

        except Exception as exc:
            db.rollback()
            doc.status = DocumentStatus.FAILED
            doc.error_message = str(exc)[:500]
            db.commit()
            logger.error("document_processing_failed", document_id=document_id, error=str(exc))
            return {"status": "error", "message": str(exc)}


ENTITY_BATCH_SIZE = 5  # chunks per extraction call


@celery_app.task(name="extract_entities", bind=True, max_retries=3)  # type: ignore[untyped-decorator]
def extract_entities(self: object, document_id: str) -> dict[str, str]:
    """Extract KG entities from an indexed document's chunks via Claude."""
    from sqlalchemy import create_engine, func
    from sqlalchemy.orm import sessionmaker

    from app.config import settings
    from app.core.brain.kg.entity_extractor import entity_extractor
    from app.models.document import Document, DocumentChunk
    from app.models.kg_entity import (
        KGEntity,
        KGEntityStatus,
        KGEntityType,
        KGEvidence,
        KGRelation,
        KGRelationType,
    )
    from app.models.kg_event import KGEvent, KGEventType

    sync_url = (
        settings.database_url
        .replace("+asyncpg", "+psycopg2")
        .replace("postgresql+psycopg2", "postgresql")
    )
    engine = create_engine(sync_url)
    session_factory = sessionmaker(bind=engine)

    with session_factory() as db:
        doc = db.execute(
            select(Document).where(Document.id == uuid.UUID(document_id))
        ).scalar_one_or_none()
        if doc is None:
            logger.error("document_not_found_for_extraction", document_id=document_id)
            return {"status": "error", "message": "Document not found"}

        try:
            # Load chunks
            chunks = db.execute(
                select(DocumentChunk)
                .where(DocumentChunk.document_id == doc.id)
                .order_by(DocumentChunk.chunk_index)
            ).scalars().all()

            if not chunks:
                logger.warning("no_chunks_for_extraction", document_id=document_id)
                return {"status": "skipped", "message": "No chunks to extract from"}

            # Get venture_id
            from app.models.venture import Venture

            venture = db.execute(
                select(Venture).where(Venture.workspace_id == doc.workspace_id)
            ).scalar_one_or_none()
            if venture is None:
                return {"status": "error", "message": "No venture found"}

            # Batch chunks and extract entities
            total_created = 0
            chunk_list = list(chunks)
            for i in range(0, len(chunk_list), ENTITY_BATCH_SIZE):
                batch = chunk_list[i : i + ENTITY_BATCH_SIZE]
                combined_text = "\n\n---\n\n".join(c.content for c in batch)

                extracted = entity_extractor.extract_from_text(combined_text)

                for raw_entity in extracted:
                    entity_type_str = raw_entity.get("type", "")
                    try:
                        entity_type = KGEntityType(entity_type_str)
                    except ValueError:
                        logger.warning(
                            "unknown_entity_type",
                            entity_type=entity_type_str,
                        )
                        continue

                    data = raw_entity.get("data", {})
                    if not isinstance(data, dict):
                        continue

                    confidence = float(raw_entity.get("confidence", 0.5))

                    # Check max entities per type
                    count = db.execute(
                        select(func.count())
                        .select_from(KGEntity)
                        .where(
                            KGEntity.venture_id == venture.id,
                            KGEntity.type == entity_type,
                        )
                    ).scalar_one()
                    if count >= 50:
                        continue

                    # Auto-status based on confidence
                    if confidence >= 0.85:
                        status = KGEntityStatus.CONFIRMED
                    elif confidence >= 0.60:
                        status = KGEntityStatus.NEEDS_REVIEW
                    else:
                        status = KGEntityStatus.SUGGESTED

                    # Check for conflict
                    conflict = None
                    name = data.get("name")
                    if name and isinstance(name, str):
                        from sqlalchemy import String, cast

                        conflict = db.execute(
                            select(KGEntity).where(
                                KGEntity.venture_id == venture.id,
                                KGEntity.type == entity_type,
                                cast(KGEntity.data, String).ilike(f"%{name}%"),
                            ).limit(1)
                        ).scalar_one_or_none()

                    entity = KGEntity(
                        venture_id=venture.id,
                        type=entity_type,
                        status=KGEntityStatus.NEEDS_REVIEW if conflict else status,
                        data=data,
                        confidence=confidence,
                    )
                    db.add(entity)
                    db.flush()

                    # Evidence link
                    evidence = KGEvidence(
                        entity_id=entity.id,
                        snippet=combined_text[:500],
                        document_id=doc.id,
                        source_type="document",
                        agent_id="entity_extractor",
                    )
                    db.add(evidence)

                    # Event
                    event = KGEvent(
                        venture_id=venture.id,
                        event_type=KGEventType.ENTITY_CREATED,
                        entity_id=str(entity.id),
                        payload={
                            "type": entity_type.value,
                            "data": data,
                            "confidence": confidence,
                        },
                        actor="system:entity_extractor",
                    )
                    db.add(event)

                    # Handle conflict
                    if conflict:
                        relation = KGRelation(
                            from_entity_id=entity.id,
                            to_entity_id=conflict.id,
                            type=KGRelationType.CONFLICTS_WITH,
                        )
                        db.add(relation)
                        if conflict.status != KGEntityStatus.PINNED:
                            conflict.status = KGEntityStatus.NEEDS_REVIEW

                        conflict_event = KGEvent(
                            venture_id=venture.id,
                            event_type=KGEventType.CONFLICT_DETECTED,
                            entity_id=str(entity.id),
                            payload={"conflicting_entity_id": str(conflict.id)},
                            actor="system:entity_extractor",
                        )
                        db.add(conflict_event)

                    total_created += 1

            db.commit()
            logger.info(
                "entity_extraction_complete",
                document_id=document_id,
                entities_created=total_created,
            )
            return {"status": "success", "entities_created": str(total_created)}

        except Exception as exc:
            db.rollback()
            logger.error(
                "entity_extraction_failed",
                document_id=document_id,
                error=str(exc),
            )
            return {"status": "error", "message": str(exc)}
