"""Tests for PptxExporter and presentation tool registration."""

import io
from typing import Any

import pytest
from pptx import Presentation

from app.core.artifacts.exporters.pptx_exporter import (
    PptxExporter,
    _detect_layout,
    pptx_exporter,
)
from app.models.artifact import ArtifactType

# ------------------------------------------------------------------ #
# Layout detection
# ------------------------------------------------------------------ #


@pytest.mark.parametrize(
    ("title", "index", "total", "expected"),
    [
        ("Welcome to Our Startup", 0, 10, "title"),
        ("Problem", 1, 10, "content"),
        ("Solution", 2, 10, "content"),
        ("Market Overview", 3, 10, "section"),
        ("Traction", 4, 10, "content"),
        ("Before vs After", 5, 10, "two_column"),
        ("The Ask", 9, 10, "closing"),
        ("Thank You", 9, 10, "closing"),
        ("Next Steps", 9, 10, "closing"),
        ("Agenda", 1, 10, "section"),
        ("Comparison", 3, 10, "two_column"),
    ],
)
def test_detect_layout(
    title: str, index: int, total: int, expected: str,
) -> None:
    assert _detect_layout(title, index, total) == expected


# ------------------------------------------------------------------ #
# Exporter basics
# ------------------------------------------------------------------ #


def _sample_content(slide_count: int = 3) -> dict[str, Any]:
    """Build a minimal valid DeckOutlineContent dict."""
    slides = []
    titles = ["Cover Slide", "Problem", "The Ask"]
    for i in range(slide_count):
        slides.append({
            "title": titles[i] if i < len(titles) else f"Slide {i + 1}",
            "key_points": [f"Point {j}" for j in range(1, 4)],
            "visual_suggestion": "Chart showing growth",
            "speaker_notes": f"Notes for slide {i + 1}",
        })
    return {"slides": slides}


def test_pptx_exporter_returns_bytes() -> None:
    """Exporter should return valid bytes that can be opened."""
    content = _sample_content()
    result = pptx_exporter.export(
        ArtifactType.DECK_OUTLINE, "Test Deck", content,
    )
    assert isinstance(result, bytes)
    assert len(result) > 0

    # Verify it's a valid PPTX by opening with python-pptx
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 3


def test_pptx_exporter_slide_count() -> None:
    """Number of generated slides should match input."""
    for count in [1, 5, 10]:
        titles = [f"Slide {i}" for i in range(count)]
        content = {
            "slides": [
                {"title": t, "key_points": ["A", "B"]}
                for t in titles
            ],
        }
        result = pptx_exporter.export(
            ArtifactType.DECK_OUTLINE, "Deck", content,
        )
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == count


def test_pptx_exporter_speaker_notes() -> None:
    """Speaker notes should be attached to each slide."""
    content = _sample_content(2)
    result = pptx_exporter.export(
        ArtifactType.DECK_OUTLINE, "Notes Test", content,
    )
    prs = Presentation(io.BytesIO(result))
    for i, slide in enumerate(prs.slides):
        notes = slide.notes_slide.notes_text_frame.text
        assert f"Notes for slide {i + 1}" in notes


def test_pptx_exporter_rejects_non_deck() -> None:
    """Should raise ValueError for non-deck_outline types."""
    with pytest.raises(ValueError, match="only supports deck_outline"):
        pptx_exporter.export(
            ArtifactType.LEAN_CANVAS, "Test", {"slides": []},
        )


def test_pptx_exporter_rejects_empty_slides() -> None:
    """Should raise ValueError when slides list is empty."""
    with pytest.raises(ValueError, match="No slides"):
        pptx_exporter.export(
            ArtifactType.DECK_OUTLINE, "Empty", {"slides": []},
        )


# ------------------------------------------------------------------ #
# Singleton
# ------------------------------------------------------------------ #


def test_pptx_exporter_singleton() -> None:
    """Module-level pptx_exporter should be a PptxExporter instance."""
    assert isinstance(pptx_exporter, PptxExporter)


# ------------------------------------------------------------------ #
# Tool registration
# ------------------------------------------------------------------ #


def test_presentation_tool_registration() -> None:
    """generate_presentation should be in the tool registry."""
    from app.core.tools.presentation_tools import (
        register_presentation_tools,
    )
    from app.core.tools.registry import tool_registry

    register_presentation_tools()
    assert tool_registry.get_definition("generate_presentation") is not None


def test_deck_architect_has_generate_presentation() -> None:
    """deck-architect's tool list should include generate_presentation."""
    from app.core.tools.registry import AGENT_TOOL_MAP

    assert "generate_presentation" in AGENT_TOOL_MAP["deck-architect"]
